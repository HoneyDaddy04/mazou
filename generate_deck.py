"""
Generate the definitive Mazou pitch deck.
Dark theme, green accents, clean typography.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──
BG_DARK = RGBColor(0x14, 0x14, 0x1E)
BG_CARD = RGBColor(0x1E, 0x1E, 0x2E)
BG_CARD_LIGHT = RGBColor(0x28, 0x28, 0x3A)
GREEN = RGBColor(0x00, 0xC8, 0x96)
GREEN_DIM = RGBColor(0x00, 0xA0, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD8)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
RED = RGBColor(0xFF, 0x45, 0x45)
TEAL_TEXT = RGBColor(0x00, 0xC8, 0x96)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox, tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = GREEN
        shape.line.width = Pt(1)
    # Round corners
    shape.adjustments[0] = 0.05
    return shape


def add_accent_line(slide, left, top, width, color=GREEN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

# Left green accent bar
add_rect(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, GREEN)

# Green top line
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GREEN)

# Logo text
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(8), Inches(1.5),
             "mazou.", font_size=72, color=WHITE, bold=True,
             font_name="Georgia")

# Green dot on the period - we'll add it as a separate element
add_text_box(slide, Inches(1.2), Inches(3.3), Inches(11), Inches(0.8),
             "The AI infrastructure layer for emerging markets",
             font_size=32, color=WHITE, bold=False, font_name="Calibri")

# Accent line
add_accent_line(slide, Inches(1.2), Inches(4.3), Inches(2.5), GREEN)

# Subtitle
add_text_box(slide, Inches(1.2), Inches(4.8), Inches(10), Inches(0.5),
             "One API. Every model. Pay in Local Currency. Waste nothing.",
             font_size=20, color=GRAY, font_name="Calibri")

# Date
add_text_box(slide, Inches(1.2), Inches(6.2), Inches(4), Inches(0.4),
             "March 2026", font_size=14, color=GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: I LIVED THIS PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GREEN)

# Title
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
             "I lived this problem.", font_size=48, color=RGBColor(0x1A, 0x1A, 0x2E),
             bold=True, font_name="Georgia")

# Left accent bar for the quote block
add_rect(slide, Inches(0.8), Inches(1.8), Pt(5), Inches(4.8), GREEN)

# Story text
txBox, tf = add_text_box(
    slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(5),
    "I lead AI Operations at a YC-backed company in Nigeria. We've been rolling out AI agents across teams and pushing AI into real day-to-day workflows.",
    font_size=19, color=RGBColor(0x33, 0x33, 0x44), font_name="Calibri")

add_paragraph(tf, "", font_size=10, color=RGBColor(0x33, 0x33, 0x44))

add_paragraph(
    tf, "At first it feels simple. Ship an agent. Connect a model. Solve a problem.",
    font_size=19, color=RGBColor(0x33, 0x33, 0x44), font_name="Calibri")

add_paragraph(tf, "", font_size=10, color=RGBColor(0x33, 0x33, 0x44))

add_paragraph(
    tf, "Then the bills start coming in.",
    font_size=19, color=RGBColor(0x1A, 0x1A, 0x2E), bold=True, font_name="Calibri")

add_paragraph(tf, "", font_size=10, color=RGBColor(0x33, 0x33, 0x44))

add_paragraph(
    tf, "Different models. Different dashboards. USD charges. No clear answer to a basic question \u2014 which agent activity is costing us money and why?",
    font_size=19, color=RGBColor(0x33, 0x33, 0x44), font_name="Calibri")

add_paragraph(tf, "", font_size=10, color=RGBColor(0x33, 0x33, 0x44))

add_paragraph(
    tf, "AI spend is the new cloud spend. It starts small, grows fast, and one day the CTO looks at the bill and realises most of it was waste. We saw this with cloud costs when the Naira collapsed. AI is next \u2014 on a faster timeline.",
    font_size=19, color=RGBColor(0x33, 0x33, 0x44), font_name="Calibri")

add_paragraph(tf, "", font_size=10, color=RGBColor(0x33, 0x33, 0x44))

add_paragraph(
    tf, "If we're already feeling this at our scale, what happens when every company is running dozens of agents?",
    font_size=20, color=GREEN_DIM, bold=True, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "The Problem", font_size=48, color=WHITE, bold=True,
             font_name="Georgia")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "African companies are adopting AI fast. But the infrastructure is broken.",
             font_size=18, color=GRAY, font_name="Calibri")

# 4 problem cards in 2x2 grid
problems = [
    ("\U0001F511  10 API keys, 10 dashboards",
     "Every model provider wants its own account, billing, and portal. Nobody has the full picture."),
    ("\U0001F4B5  USD-only billing",
     "Getting a USD virtual card, managing FX, justifying dollar charges to finance \u2014 it's a tax on AI adoption in Africa."),
    ("\U0001F50D  Zero visibility",
     "Which feature burns tokens? Which agent is expensive? Which model is overkill for the task? Nobody knows."),
    ("\U0001F30D  African models are invisible",
     "Models for Yoruba, Hausa, Swahili exist but they're scattered and impossible to access alongside GPT or Claude."),
]

card_w = Inches(5.5)
card_h = Inches(2.2)
x_positions = [Inches(0.8), Inches(6.8)]
y_positions = [Inches(2.0), Inches(4.5)]

for i, (title, desc) in enumerate(problems):
    x = x_positions[i % 2]
    y = y_positions[i // 2]
    add_rect(slide, x, y, card_w, card_h, BG_CARD)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(1.0),
                 desc, font_size=15, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: WHY NOW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "Why Now", font_size=48, color=WHITE, bold=True, font_name="Georgia")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Three forces are converging to create this opportunity right now.",
             font_size=18, color=GRAY)

why_now = [
    ("\u26A1 AI agents are going mainstream",
     "Organizations are deploying fleets of AI agents across products. The need for granular cost visibility \u2014 per agent, per activity, per task \u2014 is about to explode. No one in emerging markets is solving this.",
     GREEN),
    ("\U0001F310 Africa's developer population is surging",
     "2.5M+ developers, growing 3.5x faster than the global average. They're building with AI today, forced through painful USD workarounds. The demand is proven \u2014 the infrastructure isn't.",
     GREEN),
    ("\U0001F9E0 African AI models are emerging \u2014 with no home",
     "Models optimized for Yoruba, Hausa, Igbo, Swahili are being built. But there's no unified gateway to access them alongside GPT or Claude. First mover to aggregate wins the catalogue.",
     GREEN),
]

for i, (title, desc, accent) in enumerate(why_now):
    y = Inches(2.0) + Inches(i * 1.7)
    # Card background
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.4), BG_CARD)
    # Left accent bar
    add_rect(slide, Inches(0.8), y, Pt(5), Inches(1.4), accent)
    # Title
    add_text_box(slide, Inches(1.3), y + Inches(0.15), Inches(10.5), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    # Description
    add_text_box(slide, Inches(1.3), y + Inches(0.6), Inches(10.5), Inches(0.7),
                 desc, font_size=14, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: THE SOLUTION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "The Solution", font_size=48, color=WHITE, bold=True, font_name="Georgia")

txBox, tf = add_text_box(
    slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8),
    "A unified AI API gateway where companies across emerging markets access every major LLM through",
    font_size=18, color=GRAY)
add_paragraph(tf, "one intelligent API, billed in local currency, with full cost intelligence.",
              font_size=18, color=GREEN, bold=True)

solutions = [
    ("Local Currency Billing", "Naira, Cedi, Rand, Shilling \u2014 any supported currency. No USD cards needed."),
    ("Unified API", "One integration for OpenAI, Anthropic, Google, Mistral, DeepSeek + African AI models."),
    ("Cost Intelligence", "Per agent, per activity, per product. Know exactly what consumes your tokens and why."),
    ("Smart Routing", "Route by language, budget, complexity. African languages to African models automatically."),
    ("Budget Controls", "Set limits per agent, product, or team. Automatic throttling at thresholds."),
    ("African Model Catalogue", "First gateway to unify 9+ African LLMs + global models under one API."),
]

card_w = Inches(3.7)
card_h = Inches(2.0)

for i, (title, desc) in enumerate(solutions):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (card_w + Inches(0.2))
    y = Inches(2.5) + row * (card_h + Inches(0.2))
    add_rect(slide, x, y, card_w, card_h, BG_CARD)
    add_accent_line(slide, x, y, card_w, GREEN)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.3), card_w - Inches(0.5), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.9),
                 desc, font_size=13, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: HOW IT WORKS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "How It Works", font_size=48, color=WHITE, bold=True, font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Drop-in replacement. Five minutes to integrate.",
             font_size=18, color=GRAY)

# Code block background
code_rect = add_rect(slide, Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.5),
                     RGBColor(0x0D, 0x0D, 0x14))

# Code filename header
add_rect(slide, Inches(0.8), Inches(2.0), Inches(6.5), Inches(0.4),
         RGBColor(0x1A, 0x1A, 0x28))
add_text_box(slide, Inches(1.0), Inches(2.05), Inches(4), Inches(0.35),
             "mazou_example.py", font_size=12, color=GRAY, font_name="Consolas")

# Code content
code = """from mazou import Mazou

client = Mazou(api_key="mz_live_...")

response = client.chat.completions.create(
    model="auto",
    budget="low",
    tag="support-bot",
    messages=[{
        "role": "user",
        "content": "Bawo ni mo se le..."
    }]
)"""

add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5.8), Inches(3.8),
             code, font_size=14, color=GREEN, font_name="Consolas")

# Right side: numbered steps
steps = [
    ("1", "Single API call", "Model, budget, and tag params"),
    ("2", "Language detected", "Yoruba \u2192 African language model"),
    ("3", "Budget-aware routing", "Cheapest model meeting quality bar"),
    ("4", "Tracked & billed locally", "Tagged, billed in Naira, visible in dashboard"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.2) + Inches(i * 1.1)
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.0), y, Inches(0.45), Inches(0.45)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = num
    tf_c.paragraphs[0].font.size = Pt(16)
    tf_c.paragraphs[0].font.color.rgb = WHITE
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.word_wrap = False

    add_text_box(slide, Inches(8.7), y - Inches(0.05), Inches(4), Inches(0.35),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.7), y + Inches(0.3), Inches(4), Inches(0.35),
                 desc, font_size=13, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: WHAT WE'VE BUILT (NEW - shows traction)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "What We've Built", font_size=48, color=WHITE, bold=True, font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Not a slide deck. A working product.",
             font_size=18, color=GRAY)

built_items = [
    ("30 API endpoints", "Full REST API \u2014 auth, gateway, billing, routing, dashboard, webhooks"),
    ("40+ models priced", "Real-time cost calculation in kobo for every major global + African model"),
    ("Smart routing engine", "Language detection, budget-aware routing, complexity-based model selection"),
    ("Naira billing via Paystack", "Atomic wallet debits, auto-invoicing, full transaction history"),
    ("Cost intelligence dashboard", "Per-agent, per-model, per-feature spend tracking with recommendations"),
    ("OpenAI-compatible API", "Drop-in replacement \u2014 just change the base URL. Works with existing SDKs"),
    ("BYOK + Managed modes", "Both paths live \u2014 bring your keys or buy credits in Naira"),
    ("9+ African LLMs catalogued", "InkubaLM, YarnGPT, AfroLM, Vulavula, Khaya, Lesan, Sunflower & more"),
]

for i, (title, desc) in enumerate(built_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(2.0) + row * Inches(1.2)
    # Green bullet
    add_text_box(slide, x, y + Inches(0.05), Inches(0.3), Inches(0.3),
                 "\u2713", font_size=16, color=GREEN, bold=True)
    add_text_box(slide, x + Inches(0.35), y, Inches(5.2), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.35), y + Inches(0.4), Inches(5.2), Inches(0.35),
                 desc, font_size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: AFRICAN MODEL CATALOGUE (Differentiator)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "African Model Catalogue", font_size=48, color=WHITE, bold=True,
             font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "The first gateway to unify African + global AI models under one API.",
             font_size=18, color=GRAY)

models = [
    ("InkubaLM", "Lelapa AI", "Multilingual LLM \u2014 Swahili, Yoruba, isiXhosa, Hausa, isiZulu"),
    ("YarnGPT / SabiYarn", "Nigeria", "Yoruba, Hausa, Igbo, Pidgin text generation"),
    ("AfroLM", "Masakhane", "23 African languages \u2014 NER, sentiment, classification"),
    ("Vulavula", "Lelapa AI", "NLP for Global South \u2014 isiZulu, Sesotho, Afrikaans"),
    ("Khaya AI", "GhanaNLP", "Translation & ASR \u2014 Twi, Ewe, Ga, Dagbani"),
    ("Lesan AI", "Ethiopia", "Translation \u2014 Tigrinya, Amharic"),
    ("Sunflower", "Sunbird AI", "Open-source LLM for 31 Ugandan languages"),
    ("AfriBERTa", "Masakhane", "Embeddings for 11 African languages"),
    ("UlizaLlama", "Jacaranda", "7B LLM \u2014 321M Swahili tokens, maternal health"),
]

# Table-like layout
# Header
add_rect(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5), GREEN)
add_text_box(slide, Inches(1.0), Inches(2.05), Inches(2.5), Inches(0.4),
             "Model", font_size=14, color=WHITE, bold=True)
add_text_box(slide, Inches(3.8), Inches(2.05), Inches(2), Inches(0.4),
             "Provider", font_size=14, color=WHITE, bold=True)
add_text_box(slide, Inches(6.0), Inches(2.05), Inches(6), Inches(0.4),
             "Capabilities", font_size=14, color=WHITE, bold=True)

for i, (model, provider, caps) in enumerate(models):
    y = Inches(2.55) + Inches(i * 0.5)
    bg_color = BG_CARD if i % 2 == 0 else BG_DARK
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.48), bg_color)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 model, font_size=13, color=GREEN, bold=True)
    add_text_box(slide, Inches(3.8), y + Inches(0.05), Inches(2), Inches(0.35),
                 provider, font_size=13, color=LIGHT_GRAY)
    add_text_box(slide, Inches(6.0), y + Inches(0.05), Inches(6), Inches(0.35),
                 caps, font_size=13, color=GRAY)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "One-click activation \u2014 browse the catalogue, activate a model, immediately available through your API.",
             font_size=14, color=GREEN, bold=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "Market Opportunity", font_size=48, color=WHITE, bold=True, font_name="Georgia")

# Three big stat cards
stats = [
    ("$2.6B", "African AI market by 2030", "Growing 35% CAGR"),
    ("3,400+", "Tech companies in Nigeria alone", "All adopting AI tools"),
    ("$0", "AI infrastructure for Africa", "Nobody has built this yet"),
]

for i, (number, label, sub) in enumerate(stats):
    x = Inches(0.8) + Inches(i * 4.1)
    add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(2.8), BG_CARD)
    add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(3.1), Inches(0.8),
                 number, font_size=52, color=GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")
    add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(3.1), Inches(0.5),
                 label, font_size=15, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(3.1), Inches(0.4),
                 sub, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom context
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
             "First movers: Fintechs (Paystack, Flutterwave ecosystem), neobanks, logistics, and SaaS companies",
             font_size=15, color=GRAY)
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.5),
             "already spending $500\u2013$5,000/month on AI APIs in USD.",
             font_size=15, color=GRAY)

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
             "Global comps: Portkey ($49M raised), Helicone, LiteLLM \u2014 none solve for Africa.",
             font_size=15, color=YELLOW, bold=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "Competitive Landscape", font_size=48, color=WHITE, bold=True,
             font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
             "Global gateways exist. None serve emerging markets.",
             font_size=18, color=GRAY)

# Table
headers = ["", "OpenRouter", "LiteLLM", "Portkey", "Kong", "Mazou"]
rows = [
    ["Local Currency Billing", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["African AI Models", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["Per-Agent Cost Intel", "\u2717", "Limited", "Limited", "\u2717", "\u2713"],
    ["Language-Aware Routing", "\u2717", "\u2717", "\u2717", "\u2717", "\u2713"],
    ["Multi-Provider Gateway", "\u2713", "\u2713", "\u2713", "\u2713", "\u2713"],
    ["Smart Routing", "\u2713", "\u2713", "Limited", "\u2717", "Roadmap"],
    ["Governance & Compliance", "\u2717", "\u2717", "\u2713", "\u2713", "Roadmap"],
]

col_widths = [Inches(2.5), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(2.0)]
x_start = Inches(0.8)

# Header row
hdr_y = Inches(1.8)
add_rect(slide, x_start, hdr_y, Inches(11.3), Inches(0.45), BG_CARD_LIGHT)
cx = x_start
for j, h in enumerate(headers):
    color = GREEN if h == "Mazou" else WHITE
    add_text_box(slide, cx + Inches(0.1), hdr_y + Inches(0.05), col_widths[j] - Inches(0.2),
                 Inches(0.35), h, font_size=13, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    cx += col_widths[j]

# Data rows
for i, row in enumerate(rows):
    ry = Inches(2.3) + Inches(i * 0.55)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_rect(slide, x_start, ry, Inches(11.3), Inches(0.5), bg)
    cx = x_start
    for j, cell in enumerate(row):
        if cell == "\u2713":
            color = GREEN
        elif cell == "\u2717":
            color = RED
        elif cell == "Limited":
            color = YELLOW
        elif cell == "Roadmap":
            color = GRAY
        else:
            color = LIGHT_GRAY
        align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        add_text_box(slide, cx + Inches(0.1), ry + Inches(0.07),
                     col_widths[j] - Inches(0.2), Inches(0.35),
                     cell, font_size=12, color=color,
                     bold=(cell == "\u2713" or cell == "\u2717"),
                     alignment=align)
        cx += col_widths[j]

# Bottom insight
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.7),
             "Mazou's four differentiators \u2014 local billing, African models, per-agent cost intelligence, language-aware routing \u2014 are structurally impossible for US/EU-focused gateways to prioritize. They have no incentive to build for this market.",
             font_size=13, color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "Business Model", font_size=48, color=WHITE, bold=True, font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Two revenue streams. Both high-margin.",
             font_size=18, color=GRAY)

# Managed card
add_rect(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.2), BG_CARD)
add_accent_line(slide, Inches(0.8), Inches(2.0), Inches(5.5), GREEN)
add_text_box(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(0.4),
             "M A N A G E D", font_size=13, color=GREEN, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.7), Inches(5), Inches(0.5),
             '"Buy Through Us"', font_size=24, color=WHITE, bold=True, font_name="Georgia")
txBox, tf = add_text_box(slide, Inches(1.1), Inches(3.4), Inches(5), Inches(0.8),
                         "Customer buys Naira credits.", font_size=15, color=GRAY)
add_paragraph(tf, "We hold the API keys.", font_size=15, color=GRAY)
add_paragraph(tf, "One invoice. Zero USD.", font_size=15, color=GRAY)
add_text_box(slide, Inches(1.1), Inches(4.5), Inches(5), Inches(0.4),
             "Revenue: 10\u201320% margin on API credits", font_size=15, color=GREEN, bold=True)

# BYOK card
add_rect(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.2), BG_CARD)
add_accent_line(slide, Inches(6.8), Inches(2.0), Inches(5.5), YELLOW)
add_text_box(slide, Inches(7.1), Inches(2.3), Inches(5), Inches(0.4),
             "B Y O K", font_size=13, color=YELLOW, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.7), Inches(5), Inches(0.5),
             '"Bring Your Own Keys"', font_size=24, color=WHITE, bold=True, font_name="Georgia")
txBox, tf = add_text_box(slide, Inches(7.1), Inches(3.4), Inches(5), Inches(0.8),
                         "Customer keeps existing contracts.", font_size=15, color=GRAY)
add_paragraph(tf, "Plugs in their API keys.", font_size=15, color=GRAY)
add_paragraph(tf, "Gets the control layer on top.", font_size=15, color=GRAY)
add_text_box(slide, Inches(7.1), Inches(4.5), Inches(5), Inches(0.4),
             "Revenue: Platform subscription", font_size=15, color=YELLOW, bold=True)

# Path to scale
add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.3), BG_CARD)
add_accent_line(slide, Inches(0.8), Inches(5.6), Inches(11.5), GREEN)
add_text_box(slide, Inches(1.1), Inches(5.8), Inches(3), Inches(0.4),
             "Path to Scale", font_size=16, color=GREEN, bold=True)
add_text_box(slide, Inches(1.1), Inches(6.2), Inches(11), Inches(0.6),
             "Managed mode acquires users (lower barrier), BYOK captures enterprise (higher LTV). Both modes feed the cost intelligence flywheel \u2014 more data = better optimization recommendations = stickier product.",
             font_size=13, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: THE VISION (NEW - from highlights doc)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "The Bigger Picture", font_size=48, color=WHITE, bold=True, font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "We're not just an API gateway. We're becoming the cost control plane for the entire AI stack.",
             font_size=18, color=GRAY)

# Central insight
add_rect(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2), BG_CARD, border=True)
add_text_box(slide, Inches(1.3), Inches(2.2), Inches(10.5), Inches(0.8),
             "As AI agents start needing wallets, web search, commerce \u2014 those services plug into the same Mazou layer. Same billing, same visibility. We become the layer you can't rip out because we have your complete spend picture.",
             font_size=16, color=WHITE)

# Future services
future = [
    ("\U0001F4B3 Agent Wallets", "Agents get crypto wallets (USDT) with spending limits, auto-top-up, and transaction tracking via Coinbase"),
    ("\U0001F6D2 Agent Commerce", "Agents purchase external services through Stripe-connected accounts, with full cost visibility"),
    ("\U0001F50E Agent Web Search", "Agents use search APIs with every query cost-tracked and attributed"),
    ("\U0001F310 Content Access", "Agents browse the web via Cloudflare infrastructure, with per-fetch cost tracking"),
]

for i, (title, desc) in enumerate(future):
    x = Inches(0.8) + (i % 2) * Inches(6.0)
    y = Inches(3.6) + (i // 2) * Inches(1.5)
    add_rect(slide, x, y, Inches(5.5), Inches(1.2), BG_CARD)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.9), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(4.9), Inches(0.55),
                 desc, font_size=12, color=GRAY)

# Bottom line
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "Unified service billing \u2014 LLM costs, agent wallet spend, search costs, commerce purchases all in one dashboard.",
             font_size=14, color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: ROADMAP
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "Roadmap", font_size=48, color=WHITE, bold=True, font_name="Georgia")

phases = [
    ("Q1 2026", "Build & Beta",
     ["Core gateway live", "3 providers (OpenAI, Anthropic, Awarri)", "Naira billing via Paystack",
      "Dashboard shipped", "5 beta customers"], GREEN),
    ("Q2 2026", "Public Launch",
     ["Public launch", "8+ models integrated", "Smart routing live",
      "BYOK mode shipped", "25 paying customers"], RGBColor(0x00, 0xA0, 0x78)),
    ("Q3\u2013Q4 2026", "Scale",
     ["Full African catalogue", "ML-based routing", "Response caching",
      "Kenya + South Africa", "100 customers"], RGBColor(0x00, 0x80, 0x60)),
    ("2027", "Expand",
     ["Enterprise tier & SOC 2", "Agent wallets & commerce", "Agent monitoring",
      "Partner integrations", "250+ customers"], RGBColor(0x00, 0x60, 0x48)),
]

card_w = Inches(2.7)
for i, (timeline, title, items, accent) in enumerate(phases):
    x = Inches(0.8) + Inches(i * 3.1)
    add_rect(slide, x, Inches(1.5), card_w, Inches(5.2), BG_CARD)
    add_accent_line(slide, x, Inches(1.5), card_w, accent)

    add_text_box(slide, x + Inches(0.2), Inches(1.75), card_w - Inches(0.4), Inches(0.3),
                 timeline, font_size=13, color=accent, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.1), card_w - Inches(0.4), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Georgia")

    txBox, tf = add_text_box(slide, x + Inches(0.3), Inches(2.9),
                              card_w - Inches(0.6), Inches(3.5),
                              "", font_size=13, color=GRAY)
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "\u2022  " + item
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_before = Pt(6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: THE TEAM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "The Team", font_size=48, color=WHITE, bold=True, font_name="Georgia")

# Founder card
add_rect(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), BG_CARD)

add_text_box(slide, Inches(1.3), Inches(1.9), Inches(8), Inches(0.5),
             "Kabiru Awulu", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_text_box(slide, Inches(1.3), Inches(2.5), Inches(8), Inches(0.4),
             "Founder & CEO", font_size=16, color=GREEN, bold=True)

credentials = [
    "Head of AI Operations, Curacel (YC S22) \u2014 Led AI agent deployment and cost governance in production",
    "Former Head of IT & Service Delivery, Helium Health (YC S20)",
    "Former Global Service Operations & Automations Expert, Ericsson",
    "13+ years in telecom, fintech, and healthtech infrastructure",
    "Built and managed AI operations at scale \u2014 lived the exact problem Mazou solves",
]

txBox, tf = add_text_box(slide, Inches(1.3), Inches(3.2), Inches(10), Inches(2.5),
                         "", font_size=15, color=LIGHT_GRAY)
for j, cred in enumerate(credentials):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.text = "\u2022  " + cred
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Calibri"
    p.space_before = Pt(10)

# Why this founder
add_rect(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7), BG_CARD_LIGHT)
add_text_box(slide, Inches(1.3), Inches(6.35), Inches(10.5), Inches(0.55),
             "\U0001F3AF  Founder-market fit: Not solving a theoretical problem. This is a daily operational pain point I manage right now, at a YC company, shipping AI agents in Africa.",
             font_size=14, color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: THE ASK (NEW)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9),
             "The Ask", font_size=48, color=WHITE, bold=True, font_name="Georgia")

# Big number
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1),
             "[Raising $XXX]", font_size=44, color=GREEN, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Georgia")

add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.5),
             "[Pre-seed / Seed] to go from beta to public launch",
             font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Use of funds
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.4),
             "Use of Funds", font_size=18, color=WHITE, bold=True)

funds = [
    ("Engineering", "Hire 2 engineers \u2014 backend (Python) + frontend (React/Next.js)"),
    ("Infrastructure", "Cloud hosting, API provider deposits, monitoring & security"),
    ("Go-to-Market", "Beta program, developer relations, first 25 paying customers"),
    ("Operations", "Legal, compliance groundwork, 12-month runway"),
]

for i, (category, detail) in enumerate(funds):
    y = Inches(4.4) + Inches(i * 0.6)
    add_text_box(slide, Inches(2.5), y, Inches(2), Inches(0.4),
                 category, font_size=15, color=GREEN, bold=True)
    add_text_box(slide, Inches(4.8), y, Inches(6), Inches(0.4),
                 detail, font_size=15, color=GRAY)

# Note
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "\u2191  Update the amount and round type before sending. This slide is a placeholder for you to customize.",
             font_size=12, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: CLOSE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

# Left green accent bar
add_rect(slide, Inches(0), Inches(0), Pt(6), SLIDE_H, GREEN)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GREEN)

# Logo
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(8), Inches(1.5),
             "mazou.", font_size=72, color=WHITE, bold=True, font_name="Georgia")

# Tagline
add_text_box(slide, Inches(1.2), Inches(3.2), Inches(11), Inches(1),
             "One API. Every model.\nPay in Local Currency. Waste nothing.",
             font_size=28, color=WHITE, font_name="Calibri")

# Accent line
add_accent_line(slide, Inches(1.2), Inches(4.5), Inches(2.5), GREEN)

# The Paystack line
add_text_box(slide, Inches(1.2), Inches(5.0), Inches(11), Inches(0.6),
             "Paystack built the payment rails for Africa's internet economy.\nWe're building the infrastructure layer for Africa's AI economy.",
             font_size=17, color=GRAY, font_name="Calibri")

# Contact
add_text_box(slide, Inches(1.2), Inches(6.0), Inches(11), Inches(0.4),
             "team@vaigence.com  |  mazou.vaigence.com",
             font_size=15, color=GRAY)

# CTA
add_text_box(slide, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4),
             "Let's build this.", font_size=18, color=GREEN, bold=False)

# ── MEANING: bottom corner
add_text_box(slide, Inches(8.5), Inches(6.8), Inches(4), Inches(0.4),
             'mazou (Hausa) \u2014 "honour, dignity, respect"',
             font_size=11, color=RGBColor(0x66, 0x66, 0x77))


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = r"c:\Users\awulu\Documents\Main Development Projects\mazou\pitches\Mazou_Master_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
